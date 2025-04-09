from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    # Create a new presentation
    prs = Presentation()

    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Sample Presentation"
    subtitle.text = "Created with Python-PPTX"

    # Add a content slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Agenda'

    tf = body_shape.text_frame
    tf.text = 'First Bullet Point'

    # Add a second bullet point
    p = tf.add_paragraph()
    p.text = 'Second Bullet Point'
    p.level = 1  # Makes this a sub-bullet

    # Save the presentation
    prs.save('example.pptx')

if __name__ == '__main__':
    create_presentation()